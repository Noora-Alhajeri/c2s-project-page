from __future__ import annotations

import json
import re
from pathlib import Path

from pypdf import PdfReader
from pptx import Presentation


ROOT = Path(__file__).resolve().parents[2]
PDF_PATH = ROOT / "docs" / "C2S.pdf"
PPTX_PATH = ROOT / "docs" / "C2S.pptx"
TEMPLATE_PATH = ROOT / "docs" / "project_page" / "index.template.html"
OUTPUT_PATH = ROOT / "docs" / "project_page" / "index.html"
REPORT_PATH = ROOT / "docs" / "project_page" / "extraction_report.json"


def _extract_title_date_from_pdf(pdf_path: Path) -> tuple[str | None, str | None, int | None]:
    if not pdf_path.exists():
        return None, None, None
    reader = PdfReader(str(pdf_path))
    date_pattern = re.compile(r"\b\d{2}\s+[A-Za-z]{3}\s+\d{4}\b")
    for idx, page in enumerate(reader.pages, start=1):
        text = page.extract_text() or ""
        lines = [line.strip() for line in text.splitlines() if line.strip()]
        if "Initial Title" in lines:
            title_idx = lines.index("Initial Title")
            title_lines = lines[:title_idx]
            title = " ".join(title_lines[-2:]).strip() if title_lines else None
            date = next((line for line in lines if date_pattern.search(line)), None)
            return title or None, date or None, idx
    return None, None, None


def _extract_author_affiliation_from_pptx(pptx_path: Path) -> tuple[str | None, str | None, int | None]:
    if not pptx_path.exists():
        return None, None, None
    prs = Presentation(str(pptx_path))
    name_pattern = re.compile(r"^[A-Z][a-z]+(?:\s+[A-Z][a-z]+){1,2}$")
    affiliation_keywords = ("MBZUAI", "University", "Institute", "College", "School")
    excluded_tokens = {
        "Initial",
        "Proposed",
        "Architecture",
        "System",
        "Systems",
        "Cognitive",
        "Proposal",
        "Thesis",
        "Title",
        "Loop",
        "Vision",
        "Sensor",
        "Brain",
        "Model",
        "Robots",
    }
    candidate_name = None
    candidate_aff = None
    candidate_slide = None

    for i in range(min(10, len(prs.slides))):
        slide = prs.slides[i]
        for shape in slide.shapes:
            if not hasattr(shape, "text"):
                continue
            text = shape.text.strip()
            if not text:
                continue
            for line in text.splitlines():
                line = line.strip()
                if not line:
                    continue
                if candidate_aff is None and any(k in line for k in affiliation_keywords):
                    candidate_aff = line
                    candidate_slide = i + 1
                if candidate_name is None and name_pattern.match(line):
                    if any(token in line for token in excluded_tokens):
                        continue
                    candidate_name = line
                    candidate_slide = i + 1

    return candidate_name, candidate_aff, candidate_slide


def _render_template(template: str, values: dict[str, str]) -> str:
    output = template
    for key, value in values.items():
        output = output.replace(f"{{{{{key}}}}}", value)
    return output


def main() -> None:
    title, date, title_page = _extract_title_date_from_pdf(PDF_PATH)
    author, affiliation, author_slide = _extract_author_affiliation_from_pptx(PPTX_PATH)

    # Fallbacks if no data was found.
    title = title or "Transform Robot Actions from Conscious to Subconscious"
    subtitle = (
        "From unskilled to skilled: a triune, brain-inspired control system that "
        "moves from slow deliberation to fast automatic execution."
    )
    author = author or "[Author Name]"
    affiliation = affiliation or "Mohamed bin Zayed University of Artificial Intelligence (MBZUAI)"
    date = date or "09 Sep 2025"

    values = {
        "PROJECT_TITLE": title,
        "PROJECT_SUBTITLE": subtitle,
        "PROJECT_AUTHORS": author,
        "PROJECT_AFFILIATION": affiliation,
        "PROJECT_LOCATION": "Abu Dhabi, UAE",
        "PROJECT_DATE": date,
        "PROJECT_PAPER_LINK": "#",
        "PROJECT_CODE_LINK": "#",
        "PROJECT_VIDEO_LINK": "#",
        "PROJECT_ABSTRACT": (
            "Humanoid robots must acquire complex manipulation skills to operate "
            "effectively in human environments. Yet current vision-language-action "
            "systems execute every task as if it were the first time: full perception, "
            "full planning, full verification -- slow, costly, and repeated endlessly. "
            "This project targets a missing capability: the transformation of conscious "
            "effort into subconscious skill via a triune hierarchical memory architecture."
        ),
    }

    template = TEMPLATE_PATH.read_text(encoding="utf-8")
    output = _render_template(template, values)
    OUTPUT_PATH.write_text(output, encoding="utf-8")

    report = {
        "title": title,
        "title_source_page": title_page,
        "date": date,
        "author": author,
        "author_source_slide": author_slide,
        "affiliation": affiliation,
        "pdf_path": str(PDF_PATH),
        "pptx_path": str(PPTX_PATH),
    }
    REPORT_PATH.write_text(json.dumps(report, indent=2), encoding="utf-8")

    print(f"Wrote project page to {OUTPUT_PATH}")
    print(f"Wrote extraction report to {REPORT_PATH}")


if __name__ == "__main__":
    main()
